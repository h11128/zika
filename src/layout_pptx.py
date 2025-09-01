"""
PPTX layout generator for 3x3 Chinese character learning cards.
"""

from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from typing import List, Dict, Tuple
import math


class PPTXCardGenerator:
    """Generator for PPTX format learning cards (rows x cols grid)."""

    def __init__(self, page_size: str = "A4", card_size_cm: float = 6.0,
                 gap_cm: float = 0.6, margin_cm: float = 1.0,
                 layout_rows: int = 3, layout_cols: int = 3, layout_auto_fill: bool = True):
        """
        Initialize the PPTX generator.

        Args:
            page_size: Page size ("A4" or "Letter")
            card_size_cm: Size of each card in cm (used when auto_fill is False)
            gap_cm: Gap between cards in cm
            margin_cm: Page margin in cm
            layout_rows: Number of rows per page
            layout_cols: Number of columns per page
            layout_auto_fill: If True, compute max card size to fill page within margins
        """
        self.page_size = page_size
        self.gap_cm = gap_cm
        self.margin_cm = margin_cm
        self.layout_rows = max(1, int(layout_rows or 3))
        self.layout_cols = max(1, int(layout_cols or 3))
        self.layout_auto_fill = bool(layout_auto_fill)

        # Set page dimensions
        if page_size.upper() == "A4":
            self.page_width = Cm(21.0)   # A4 width
            self.page_height = Cm(29.7)  # A4 height
        elif page_size.upper() == "LETTER":
            self.page_width = Inches(8.5)   # Letter width
            self.page_height = Inches(11.0) # Letter height
        else:
            raise ValueError(f"Unsupported page size: {page_size}")

        # Page size in cm
        page_width_cm = self.page_width.cm if hasattr(self.page_width, 'cm') else self.page_width / 914400 * 2.54
        page_height_cm = self.page_height.cm if hasattr(self.page_height, 'cm') else self.page_height / 914400 * 2.54
        max_w = max(0.0, page_width_cm - 2 * self.margin_cm)
        max_h = max(0.0, page_height_cm - 2 * self.margin_cm)

        # Determine card size
        if self.layout_auto_fill:
            size_w = (max_w - max(0, self.layout_cols - 1) * self.gap_cm) / self.layout_cols if self.layout_cols > 0 else 0
            size_h = (max_h - max(0, self.layout_rows - 1) * self.gap_cm) / self.layout_rows if self.layout_rows > 0 else 0
            self.card_size_cm = max(0.0, min(size_w, size_h))
        else:
            self.card_size_cm = card_size_cm

        # Calculate grid layout and ensure it fits; shrink if needed
        self.grid_width = self.layout_cols * self.card_size_cm + max(0, self.layout_cols - 1) * self.gap_cm
        self.grid_height = self.layout_rows * self.card_size_cm + max(0, self.layout_rows - 1) * self.gap_cm

        tolerance = 0.1
        if self.grid_width + 2 * self.margin_cm > page_width_cm + tolerance or \
           self.grid_height + 2 * self.margin_cm > page_height_cm + tolerance:
            # Compute scale to fit
            scale_w = max_w / self.grid_width if self.grid_width > 0 else 1.0
            scale_h = max_h / self.grid_height if self.grid_height > 0 else 1.0
            scale = min(scale_w, scale_h) * 0.995
            self.card_size_cm *= max(0.0, scale)
            self.grid_width = self.layout_cols * self.card_size_cm + max(0, self.layout_cols - 1) * self.gap_cm
            self.grid_height = self.layout_rows * self.card_size_cm + max(0, self.layout_rows - 1) * self.gap_cm

    def create_presentation(self) -> Presentation:
        """Create a new presentation with custom page size."""
        prs = Presentation()
        
        # Set slide size
        prs.slide_width = int(self.page_width)
        prs.slide_height = int(self.page_height)
        
        return prs
    
    def add_cards_page(self, prs: Presentation, cards: List[Dict[str, str]],
                      hanzi_font_size: int = 48, pinyin_font_size: int = 18, english_font_size: int = 14,
                      hanzi_font_family: str = "Microsoft YaHei", background_color: str = "#FFFFFF") -> None:
        """
        Add a page with up to rows*cols cards in a grid.

        Args:
            prs: Presentation object
            cards: List of card dictionaries with 'hanzi', 'pinyin', 'english' keys
            hanzi_font_size: Font size for Chinese characters
            pinyin_font_size: Font size for pinyin
            english_font_size: Font size for English
        """
        # Use blank slide layout
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add up to rows*cols cards
        max_cards = self.layout_rows * self.layout_cols
        for i, card in enumerate(cards[:max_cards]):
            row = i // self.layout_cols
            col = i % self.layout_cols
            self._add_single_card(slide, card, row, col, hanzi_font_size, pinyin_font_size, english_font_size,
                                hanzi_font_family, background_color)

    def _add_single_card(self, slide, card: Dict[str, str], row: int, col: int,
                        hanzi_font_size: int, pinyin_font_size: int, english_font_size: int,
                        hanzi_font_family: str = "Microsoft YaHei", background_color: str = "#FFFFFF") -> None:
        """
        Add a single card to the slide.
        
        Args:
            slide: Slide object
            card: Card data dictionary
            row: Grid row (0-based)
            col: Grid column (0-based)
            hanzi_font_size: Font size for Chinese characters
            pinyin_font_size: Font size for pinyin
            english_font_size: Font size for English
        """
        # Calculate position (center grid within margins)
        page_width_cm = self.page_width.cm if hasattr(self.page_width, 'cm') else self.page_width / 914400 * 2.54
        page_height_cm = self.page_height.cm if hasattr(self.page_height, 'cm') else self.page_height / 914400 * 2.54
        avail_w = max(0.0, page_width_cm - 2 * self.margin_cm)
        avail_h = max(0.0, page_height_cm - 2 * self.margin_cm)
        grid_width = self.layout_cols * self.card_size_cm + max(0, self.layout_cols - 1) * self.gap_cm
        grid_height = self.layout_rows * self.card_size_cm + max(0, self.layout_rows - 1) * self.gap_cm
        start_x_cm = self.margin_cm + max(0.0, (avail_w - grid_width) / 2)
        start_y_cm = self.margin_cm + max(0.0, (avail_h - grid_height) / 2)

        left = Cm(start_x_cm + col * (self.card_size_cm + self.gap_cm))
        top = Cm(start_y_cm + row * (self.card_size_cm + self.gap_cm))
        width_cm = Cm(self.card_size_cm)
        height_cm = Cm(self.card_size_cm)

        # Create rectangle shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width_cm, height_cm
        )
        
        # Set border
        shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        shape.line.width_cm = Pt(2)
        
        # Set fill color
        shape.fill.solid()
        # Parse hex color to RGB
        if background_color.startswith('#'):
            hex_color = background_color[1:]
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
        else:
            r, g, b = 255, 255, 255  # Default to white
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        
        # Configure text frame
        text_frame = shape.text_frame
        # Avoid potential python-pptx clear() issues on some Office versions
        text_frame.text = ""
        text_frame.margin_left = Cm(0.2)
        text_frame.margin_right = Cm(0.2)
        text_frame.margin_top_px = Cm(0.2)
        text_frame.margin_bottom_px = Cm(0.2)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        text_frame.word_wrap = True
        
        # Add text content
        self._add_card_text(text_frame, card, hanzi_font_size, pinyin_font_size, english_font_size, hanzi_font_family)
    
    def _add_card_text(self, text_frame, card: Dict[str, str],
                      hanzi_font_size: int, pinyin_font_size: int, english_font_size: int,
                      hanzi_font_family: str = "Microsoft YaHei") -> None:
        """
        Add formatted text to a card's text frame.

        Args:
            text_frame: Text frame object
            card: Card data dictionary
            hanzi_font_size: Font size for Chinese characters
            pinyin_font_size: Font size for pinyin
            english_font_size: Font size for English
        """
        # Get text content
        hanzi = card.get('hanzi', '').strip()
        pinyin = card.get('pinyin', '').strip()
        english = card.get('english', '').strip()

        # Do NOT clear here; it was cleared by caller
        # Calculate spacing for even distribution
        total_elements = sum([1 for x in [hanzi, pinyin, english] if x])
        if total_elements == 0:
            return

        # Decide where to put the first line (use the default paragraph)
        para = text_frame.paragraphs[0]
        first_used = False

        # Add hanzi (Chinese characters)
        if hanzi:
            p = para if not first_used else text_frame.add_paragraph()
            first_used = True
            p.text = hanzi
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(hanzi_font_size)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            # Use custom font for Chinese characters
            try:
                p.font.name = hanzi_font_family
            except:
                try:
                    p.font.name = "Microsoft YaHei"  # Fallback
                except:
                    try:
                        p.font.name = "SimSun"
                    except:
                        p.font.name = "Arial Unicode MS"
            p.line_spacing = 1.0
            if total_elements > 1:
                p.space_after = Pt(8)

        # Add pinyin
        if pinyin:
            p = para if not first_used else text_frame.add_paragraph()
            first_used = True
            p.text = pinyin
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(pinyin_font_size)
            p.font.italic = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            # Use fonts that handle pinyin tone marks well
            try:
                p.font.name = "Calibri"  # Good for pinyin tone marks
            except:
                try:
                    p.font.name = "Arial"
                except:
                    p.font.name = "Times New Roman"
            p.line_spacing = 1.0
            if english:  # Add space only if there's English text below
                p.space_after = Pt(8)

        # Add English translation
        if english:
            p = para if not first_used else text_frame.add_paragraph()
            first_used = True
            p.text = english
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(english_font_size)
            p.font.name = "Arial"
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.line_spacing = 1.0
    
    def generate_pptx(self, cards: List[Dict[str, str]], output_path: str,
                     hanzi_font_size: int = 48, pinyin_font_size: int = 18, english_font_size: int = 14,
                     hanzi_font_family: str = "Microsoft YaHei", background_color: str = "#FFFFFF") -> bool:
        """
        Generate complete PPTX file with all cards.
        
        Args:
            cards: List of all card data
            output_path: Output file path
            hanzi_font_size: Font size for Chinese characters
            pinyin_font_size: Font size for pinyin
            english_font_size: Font size for English
            
        Returns:
            True if successful
        """
        try:
            prs = self.create_presentation()

            # Split cards into pages of rows*cols
            page_size = max(1, self.layout_rows * self.layout_cols)
            for i in range(0, len(cards), page_size):
                page_cards = cards[i:i+page_size]
                self.add_cards_page(prs, page_cards, hanzi_font_size, pinyin_font_size, english_font_size,
                                  hanzi_font_family, background_color)

            # Save presentation
            prs.save(output_path)
            return True
            
        except Exception as e:
            print(f"Error generating PPTX: {e}")
            return False
    
    def get_layout_info(self) -> Dict[str, any]:
        """Get layout information for debugging."""
        return {
            'page_size': self.page_size,
            'page_width_cm': self.page_width.cm if hasattr(self.page_width, 'cm') else self.page_width / 914400 * 2.54,
            'page_height_cm': self.page_height.cm if hasattr(self.page_height, 'cm') else self.page_height / 914400 * 2.54,
            'card_size_cm': self.card_size_cm,
            'gap_cm': self.gap_cm,
            'margin_cm': self.margin_cm,
            'layout_rows': self.layout_rows,
            'layout_cols': self.layout_cols,
            'grid_width_cm': self.grid_width,
            'grid_height_cm': self.grid_height,
            'cards_per_page': self.layout_rows * self.layout_cols
        }


def create_sample_cards() -> List[Dict[str, str]]:
    """Create sample cards for testing."""
    return [
        {'hanzi': '爱', 'pinyin': 'ài', 'english': 'love'},
        {'hanzi': '家', 'pinyin': 'jiā', 'english': 'home; family'},
        {'hanzi': '朋友', 'pinyin': 'péng yǒu', 'english': 'friend'},
        {'hanzi': '水', 'pinyin': 'shuǐ', 'english': 'water'},
        {'hanzi': '火', 'pinyin': 'huǒ', 'english': 'fire'},
        {'hanzi': '山', 'pinyin': 'shān', 'english': 'mountain'},
        {'hanzi': '月', 'pinyin': 'yuè', 'english': 'moon; month'},
        {'hanzi': '日', 'pinyin': 'rì', 'english': 'sun; day'},
        {'hanzi': '木', 'pinyin': 'mù', 'english': 'wood; tree'},
    ]


if __name__ == "__main__":
    # Test the PPTX generator
    print("Testing PPTX generator...")
    
    generator = PPTXCardGenerator(page_size="A4", card_size_cm=6.0, gap_cm=0.6, margin_cm=1.0)
    
    print("Layout info:")
    layout_info = generator.get_layout_info()
    for key, value in layout_info.items():
        print(f"  {key}: {value}")
    
    # Generate sample PPTX
    sample_cards = create_sample_cards()
    output_path = "../out/test_cards.pptx"
    
    success = generator.generate_pptx(sample_cards, output_path)
    if success:
        print(f"\nSample PPTX generated: {output_path}")
    else:
        print("\nFailed to generate PPTX")
